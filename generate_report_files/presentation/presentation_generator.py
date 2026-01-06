#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Генератор презентаций - скрипт для формирования презентаций по отчетам из БД
"""
import io
import os
import json
import psycopg2
from minio import Minio
from dotenv import load_dotenv
from typing import Dict, List, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from datetime import datetime
import requests
from io import BytesIO
from PIL import Image


# Загружаем переменные окружения
load_dotenv()


class PresentationGenerator:
    def __init__(self):
        """Инициализация подключений к БД и MinIO"""
        # Настройки БД
        self.db_config = {
            'host': os.getenv('DB_HOST', 'localhost'),
            'port': os.getenv('DB_PORT', '5432'),
            'database': os.getenv('DB_NAME'),
            'user': os.getenv('DB_USER'),
            'password': os.getenv('DB_PASSWORD')
        }
        
        # Настройки MinIO
        self.minio_client = Minio(
            endpoint=os.getenv('S3_ENDPOINT_URL', 'minio.upk-mos.ru'),
            access_key=os.getenv('S3_ACCESS_KEY'),
            secret_key=os.getenv('S3_SECRET_KEY'),
            secure=os.getenv('S3_SECURE', 'False').lower() == 'true'
        )
        
        self.bucket_name = os.getenv('S3_BUCKET_NAME', 'dit-services-dev')
        
        # Папка для результатов
        self.output_folder = 'presentations_results'
        self._ensure_output_folder()
        
        # Кэш для изображений
        self.image_cache = {}
    
    def _ensure_output_folder(self):
        """Создать папку для результатов, если её нет"""
        if not os.path.exists(self.output_folder):
            os.makedirs(self.output_folder)
            print(f"📁 Создана папка для результатов: {self.output_folder}")
        else:
            print(f"📁 Папка для результатов: {self.output_folder}")
    
    def get_pending_reports(self) -> List[Dict]:
        """Получить отчеты со статусом 1 (готовые к обработке)"""
        try:
            print(f"🔌 Подключение к БД: {self.db_config['host']}:{self.db_config['port']}/{self.db_config['database']}")
            conn = psycopg2.connect(**self.db_config)
            cursor = conn.cursor()
            
            # Устанавливаем схему по умолчанию
            cursor.execute("SET search_path TO gen_report_context_contracts, public;")
            print("✓ Схема установлена: gen_report_context_contracts")
            
            # Получаем отчеты со статусом 1 и данные заявки
            query = """
            SELECT 
                r.id, 
                r.id_contracts, 
                r.id_requests, 
                c.number_contract, 
                c.subject_contract,
                req.campany_yandex_direct
            FROM reports r
            JOIN contracts c ON r.id_contracts = c.id
            JOIN requests req ON r.id_requests = req.id
            WHERE r.id_status = 1 AND (r.is_deleted = false OR r.is_deleted IS NULL)
            ORDER BY r.create_entry DESC
            """
            
            print(f"🔍 Выполняем запрос...")
            cursor.execute(query)
            reports = []
            
            for row in cursor.fetchall():
                reports.append({
                    'id': row[0],
                    'id_contracts': row[1],
                    'id_requests': row[2],
                    'number_contract': row[3],
                    'subject_contract': row[4],
                    'campaign_ids': row[5]  # JSONB с id кампаний
                })
            
            print(f"✅ Найдено отчетов для обработки: {len(reports)}")
            
            cursor.close()
            conn.close()
            
            return reports
            
        except Exception as e:
            print(f"❌ Ошибка при получении отчетов из БД: {e}")
            return []
    
    def load_file_from_minio(self, report_id: int, filename: str) -> Optional[Dict]:
        """Загрузить JSON файл из MinIO для конкретного отчета"""
        try:
            folder_path = f"gen_report_context_contracts/data_yandex_direct/{report_id}_результаты"
            object_path = f"{folder_path}/{filename}"
            
            # Проверяем существование объекта
            if self.minio_client.stat_object(self.bucket_name, object_path):
                # Загружаем объект
                response = self.minio_client.get_object(self.bucket_name, object_path)
                content = response.read().decode('utf-8')
                data = json.loads(content)
                print(f"✓ Загружен файл {filename}")
                return data
            else:
                print(f"⚠ Файл {filename} не найден")
                return None
                
        except Exception as e:
            print(f"✗ Ошибка при загрузке {filename}: {e}")
            return None
    
    def filter_rsy_campaigns(self, campaigns_data: Dict, request_campaign_ids: Any) -> List[Dict]:
        """
        Фильтровать кампании:
        1. Только те, которые есть в заявке (request_campaign_ids)
        2. Только те, где в Name есть "РСЯ"
        """
        if not campaigns_data or not campaigns_data.get('result'):
            print("⚠ Нет данных кампаний")
            return []
        
        all_campaigns = campaigns_data.get('result', {}).get('Campaigns', [])
        print(f"📊 Всего кампаний в данных: {len(all_campaigns)}")
        
        # Получаем список ID кампаний из заявки
        allowed_campaign_ids = []
        if request_campaign_ids:
            print(f"🔍 Структура campaign_ids: {type(request_campaign_ids)}")
            print(f"🔍 Содержимое campaign_ids (первые 200 символов): {str(request_campaign_ids)[:200]}")
            
            # request_campaign_ids - это JSONB, может быть списком словарей или списком ID
            if isinstance(request_campaign_ids, list):
                for item in request_campaign_ids:
                    if isinstance(item, dict):
                        # Если это словарь, берем Id
                        campaign_id = item.get('Id') or item.get('id')
                        if campaign_id:
                            allowed_campaign_ids.append(int(campaign_id))
                    elif isinstance(item, (int, str)):
                        # Если это просто ID
                        try:
                            allowed_campaign_ids.append(int(item))
                        except ValueError:
                            pass
            elif isinstance(request_campaign_ids, dict):
                # Проверяем, есть ли ключ 'campaigns' в структуре
                if 'campaigns' in request_campaign_ids:
                    campaigns_list = request_campaign_ids['campaigns']
                    if isinstance(campaigns_list, list):
                        for item in campaigns_list:
                            if isinstance(item, dict):
                                campaign_id = item.get('Id') or item.get('id')
                                if campaign_id:
                                    allowed_campaign_ids.append(int(campaign_id))
                            elif isinstance(item, (int, str)):
                                try:
                                    allowed_campaign_ids.append(int(item))
                                except ValueError:
                                    pass
                else:
                    # Если это словарь с ID как ключами
                    for k, v in request_campaign_ids.items():
                        try:
                            allowed_campaign_ids.append(int(k))
                        except ValueError:
                            # Ключ не является числом, пропускаем
                            pass
        
        print(f"📋 ID кампаний из заявки: {allowed_campaign_ids}")
        
        # Фильтруем кампании
        rsy_campaigns = []
        for campaign in all_campaigns:
            campaign_id = campaign.get('Id')
            campaign_name = campaign.get('Name', '')
            
            # Проверяем, что кампания есть в списке из заявки
            if allowed_campaign_ids and campaign_id not in allowed_campaign_ids:
                continue
            
            # Проверяем, что в названии есть "РСЯ"
            if '/РСЯ/' in campaign_name:
                rsy_campaigns.append(campaign)
                print(f"  ✓ Найдена РСЯ-кампания: {campaign_name}")
        
        print(f"✅ Отфильтровано РСЯ-кампаний: {len(rsy_campaigns)}")
        return rsy_campaigns
    
    def extract_campaign_subtitle(self, campaign_name: str) -> str:
        """
        Извлечь подназвание кампании после "РСЯ/"
        Например: "Карта москвича/РСЯ/Брендовые ключи" -> "Брендовые ключи"
        """
        if '/РСЯ/' in campaign_name:
            parts = campaign_name.split('/РСЯ/')
            if len(parts) > 1:
                return parts[1].strip()
        return campaign_name
    
    def get_unique_ads_for_campaign(self, campaign_id: int, ads_data: Dict, image_hashes_data: Dict) -> List[Dict]:
        """
        Получить уникальные комбинации объявлений для кампании
        Уникальность определяется по (Title, Text, AdImageHash)
        """
        if not ads_data or not ads_data.get('result'):
            return []
        
        all_ads = ads_data.get('result', {}).get('Ads', [])
        
        # Фильтруем объявления по campaign_id и наличию изображения
        campaign_ads = []
        for ad in all_ads:
            if ad.get('CampaignId') == campaign_id:
                text_ad = ad.get('TextAd', {})
                ad_image_hash = text_ad.get('AdImageHash')
                
                # Только объявления с изображением
                if ad_image_hash:
                    campaign_ads.append(ad)
        
        print(f"  📊 Найдено объявлений с изображениями для кампании {campaign_id}: {len(campaign_ads)}")
        
        # Создаем уникальные комбинации
        unique_combinations = {}
        for ad in campaign_ads:
            text_ad = ad.get('TextAd', {})
            title = text_ad.get('Title', '')
            text = text_ad.get('Text', '')
            ad_image_hash = text_ad.get('AdImageHash')
            
            # Ключ для уникальности
            key = (title, text, ad_image_hash)
            
            if key not in unique_combinations:
                # Получаем URL изображения
                image_url = self._get_image_url_by_hash(ad_image_hash, image_hashes_data)
                
                # Получаем ссылку на посадочную страницу
                href = text_ad.get('Href', '')
                
                unique_combinations[key] = {
                    'title': title,
                    'text': text,
                    'image_hash': ad_image_hash,
                    'image_url': image_url,
                    'href': href
                }
        
        print(f"  ✅ Уникальных комбинаций: {len(unique_combinations)}")
        return list(unique_combinations.values())
    
    def _get_image_url_by_hash(self, ad_image_hash: str, image_hashes_data: Dict) -> Optional[str]:
        """Получить URL изображения по хешу"""
        if not image_hashes_data:
            return None
        
        ad_images = image_hashes_data.get('result', {}).get('AdImages', [])
        for image in ad_images:
            if image.get('AdImageHash') == ad_image_hash:
                return image.get('OriginalUrl')
        
        return None
    
    def download_image(self, url: str, max_retries: int = 3) -> Optional[BytesIO]:
        """Скачать изображение по URL с повторными попытками"""
        # Проверяем кэш
        if url in self.image_cache:
            print(f"    📦 Используем кэшированное изображение")
            return BytesIO(self.image_cache[url])
        
        for attempt in range(max_retries):
            try:
                print(f"    🔄 Попытка {attempt + 1}/{max_retries} скачивания изображения...")
                response = requests.get(url, timeout=15)
                if response.status_code == 200:
                    content = response.content
                    # Сохраняем в кэш
                    self.image_cache[url] = content
                    print(f"    ✅ Изображение скачано и закэшировано")
                    return BytesIO(content)
                else:
                    print(f"    ⚠ HTTP {response.status_code} на попытке {attempt + 1}")
            except Exception as e:
                print(f"    ✗ Ошибка на попытке {attempt + 1}: {e}")
                if attempt < max_retries - 1:
                    import time
                    time.sleep(2)  # Пауза перед повторной попыткой
        
        print(f"    ❌ Не удалось скачать изображение после {max_retries} попыток")
        return None
    
    def compress_image(self, image_stream: BytesIO, max_width: int = 800, max_height: int = 600, quality: int = 85) -> BytesIO:
        """
        Сжать изображение для уменьшения размера файла
        """
        try:
            image_stream.seek(0)
            img = Image.open(image_stream)
            
            # Получаем исходные размеры
            original_width, original_height = img.size
            print(f"    📏 Исходные размеры: {original_width}x{original_height}")
            
            # Вычисляем новые размеры с сохранением пропорций
            ratio = min(max_width / original_width, max_height / original_height)
            
            # Если изображение уже меньше максимальных размеров, не увеличиваем его
            if ratio >= 1:
                ratio = 1
            
            new_width = int(original_width * ratio)
            new_height = int(original_height * ratio)
            
            print(f"    📏 Новые размеры: {new_width}x{new_height} (коэффициент: {ratio:.2f})")
            
            # Изменяем размер изображения
            if ratio < 1:
                img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # Конвертируем в RGB если нужно (для JPEG)
            if img.mode in ('RGBA', 'LA', 'P'):
                # Создаем белый фон для прозрачных изображений
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                img = background
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Сохраняем с сжатием
            output = BytesIO()
            img.save(output, format='JPEG', quality=quality, optimize=True)
            output.seek(0)
            
            # Получаем размеры до и после сжатия
            original_size = len(image_stream.getvalue())
            compressed_size = len(output.getvalue())
            compression_ratio = (1 - compressed_size / original_size) * 100 if original_size > 0 else 0
            
            print(f"    💾 Сжатие: {original_size // 1024}KB → {compressed_size // 1024}KB ({compression_ratio:.1f}% экономии)")
            
            return output
            
        except Exception as e:
            print(f"    ✗ Ошибка при сжатии изображения: {e}")
            return image_stream
    
    def create_presentation(self, rsy_campaigns: List[Dict], ads_data: Dict, image_hashes_data: Dict, output_path: str) -> bool:
        """Создать презентацию с заголовками для РСЯ-кампаний"""
        try:
            # Создаем новую презентацию
            prs = Presentation()
            
            # Устанавливаем размер слайда (16:9)
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)  # 16:9 = 10/5.625
            
            # Создаем слайд для каждой РСЯ-кампании
            for campaign in rsy_campaigns:
                campaign_id = campaign.get('Id')
                campaign_name = campaign.get('Name', '')
                subtitle = self.extract_campaign_subtitle(campaign_name)
                
                # Формируем заголовок
                title_text = f"Текстово-графические объявления РСЯ: {subtitle}"
                
                print(f"  📄 Создание слайда: {title_text}")
                
                # Получаем уникальные объявления для кампании
                unique_ads = self.get_unique_ads_for_campaign(campaign_id, ads_data, image_hashes_data)
                
                if not unique_ads:
                    print(f"  ⚠ Нет объявлений с изображениями для кампании {campaign_id}")
                    continue
                
                # Определяем, нужно ли разбивать на части
                total_ads = len(unique_ads)
                if total_ads > 10:
                    # Разбиваем на две части
                    ads_per_part = total_ads // 2 + (total_ads % 2)  # Округляем вверх для первой части
                    parts = [
                        unique_ads[:ads_per_part],
                        unique_ads[ads_per_part:]
                    ]
                    print(f"  📑 Разбиваем {total_ads} объявлений на 2 части: {len(parts[0])} и {len(parts[1])} объявлений")
                else:
                    parts = [unique_ads]  # Одна часть со всеми объявлениями
                
                # Создаем слайд(ы)
                for part_idx, part_ads in enumerate(parts, 1):
                    # Формируем заголовок с учетом части
                    part_title = title_text
                    if len(parts) > 1:
                        part_title = f"{title_text} (Часть {part_idx})"
                    
                    # Добавляем пустой слайд
                    blank_slide_layout = prs.slide_layouts[6]  # 6 = пустой слайд
                    slide = prs.slides.add_slide(blank_slide_layout)
                
                    # Размеры слайда
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    
                    # === 1. ЗАГОЛОВОК (сверху) ===
                    title_top = Inches(0.2)
                    title_height = Inches(0.4)
                    
                    textbox = slide.shapes.add_textbox(Inches(0.3), title_top, Inches(9.4), title_height)
                    text_frame = textbox.text_frame
                    text_frame.word_wrap = True
                    
                    p = text_frame.paragraphs[0]
                    p.text = part_title  # Используем заголовок с номером части
                    p.alignment = PP_ALIGN.LEFT
                    
                    run = p.runs[0]
                    run.font.name = 'Proxima Nova'
                    run.font.size = Pt(16)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    
                    # === 2. КАРТИНКИ (15% высоты слайда) ===
                    images_top = title_top + title_height + Inches(0.1)
                    images_height = Inches(0.8)  # Фиксированная высота для изображений
                    
                    # Добавляем картинки для текущей части
                    self._add_images_to_slide(slide, part_ads, Inches(0.3), images_top, Inches(9.4), images_height)
                    
                    # === 3. ТАБЛИЦА (оставшееся место) ===
                    table_top = images_top + images_height + Inches(0.1)
                    table_height = slide_height - table_top - Inches(0.6)  # Оставляем место для ссылок
                    
                    # Добавляем таблицу для текущей части
                    self._add_table_to_slide(slide, part_ads, Inches(0.3), table_top, Inches(9.4), table_height)
                    
                    # === 4. БЛОК СО ССЫЛКАМИ (внизу) ===
                    links_height = Inches(0.4)
                    links_top = slide_height - Inches(0.5)  # Фиксированная позиция снизу
                    
                    # Добавляем ссылки для текущей части
                    self._add_links_textbox(slide, part_ads, Inches(0.3), links_top, Inches(9.4), links_height)
                    
                    print(f"  ✅ Создан слайд {part_idx} из {len(parts)}")
            
            # Сохраняем презентацию
            prs.save(output_path)
            print(f"✅ Презентация сохранена: {output_path}")
            return True
            
        except Exception as e:
            print(f"❌ Ошибка при создании презентации: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def _add_images_to_slide(self, slide, unique_ads: List[Dict], left, top, width, height):
        """Добавить картинки на слайд с сохранением пропорций"""
        if not unique_ads:
            return
        
        num_images = len(unique_ads)
        print(f"  🖼️ Добавление {num_images} картинок на слайд")
        
        padding = Inches(0.05)  # Отступ между картинками
        target_height = height  # Фиксированная высота для всех изображений
        
        current_left = left
        
        for i, ad in enumerate(unique_ads):
            image_url = ad.get('image_url')
            if not image_url:
                print(f"    ⚠ Нет URL изображения для объявления {i+1}")
                continue
            
            # Скачиваем изображение
            image_stream = self.download_image(image_url)
            if not image_stream:
                print(f"    ⚠ Не удалось скачать изображение {i+1}")
                continue
            
            try:
                # Сжимаем изображение перед добавлением
                compressed_stream = self.compress_image(image_stream, max_width=800, max_height=600, quality=85)
                
                # Получаем размеры сжатого изображения
                compressed_stream.seek(0)
                img = Image.open(compressed_stream)
                original_width, original_height = img.size
                
                # Вычисляем пропорциональную ширину
                aspect_ratio = original_width / original_height
                proportional_width = int(target_height * aspect_ratio)
                
                # Вычисляем доступную ширину для каждого изображения
                available_width = (width - (padding * (num_images - 1))) / num_images
                
                # Масштабируем изображение, чтобы оно поместилось в доступное пространство
                if proportional_width > available_width:
                    # Если изображение слишком широкое, уменьшаем его высоту пропорционально
                    new_height = int(target_height * (available_width / proportional_width))
                    compressed_stream.seek(0)
                    pic = slide.shapes.add_picture(compressed_stream, current_left, top + (target_height - new_height) / 2, width=available_width)
                else:
                    # Если изображение помещается, используем исходную высоту
                    compressed_stream.seek(0)
                    pic = slide.shapes.add_picture(compressed_stream, current_left, top, height=target_height)
                
                print(f"    ✅ Добавлено изображение {i+1} ({original_width}x{original_height} -> {pic.width}x{pic.height})")
                
                # Сдвигаем позицию для следующего изображения
                current_left += available_width + padding
                
            except Exception as e:
                print(f"    ✗ Ошибка добавления изображения {i+1}: {e}")
    
    def _add_table_to_slide(self, slide, unique_ads: List[Dict], left, top, width, height):
        """Добавить таблицу на слайд"""
        if not unique_ads:
            return
        
        num_rows = len(unique_ads) + 1  # +1 для заголовка
        num_cols = 2  # Заголовок и Текст
        
        print(f"  📊 Добавление таблицы {num_rows}x{num_cols}")
        
        # Создаем таблицу
        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table
        
        # Включаем первую строку (заголовок)
        table.first_row = True
        table.banded_rows = False
        table.banded_cols = False
        
        # Функция для установки черных границ для всех ячеек таблицы
        def set_table_borders():
            """Применить черные границы ко всем ячейкам таблицы"""
            for row in table.rows:
                for cell in row.cells:
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    
                    # Создаем элементы границ для всех сторон
                    for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                        # Удаляем существующую границу
                        for existing in tcPr.findall(f'.//a:{border_name}', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                            tcPr.remove(existing)
                        
                        # Создаем новую границу
                        ln = OxmlElement(f'a:{border_name}')
                        ln.set('w', '12700')  # Ширина 1pt
                        
                        # Черный цвет
                        solidFill = OxmlElement('a:solidFill')
                        srgbClr = OxmlElement('a:srgbClr')
                        srgbClr.set('val', '000000')
                        solidFill.append(srgbClr)
                        ln.append(solidFill)
                        
                        tcPr.append(ln)
        
        # Применяем границы
        set_table_borders()
        
        # Настраиваем ширину столбцов (40% для заголовка, 60% для текста)
        table.columns[0].width = int(width * 0.4)
        table.columns[1].width = int(width * 0.6)
        
        # Заголовки таблицы
        header_cells = [table.cell(0, 0), table.cell(0, 1)]
        header_texts = ['Заголовок', 'Текст']
        
        for i, (cell, text) in enumerate(zip(header_cells, header_texts)):
            cell.text = text
            # Заливка цветом #fff2cd
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 242, 205)  # #fff2cd
            
            # Форматирование текста
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Times New Roman'
                    run.font.size = Pt(12)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 0, 0)  # Черный цвет текста
        
        # Заполняем данные
        font_size = Pt(12)
        
        # Пробуем поместить все данные с начальным шрифтом
        for row_idx, ad in enumerate(unique_ads, start=1):
            title_cell = table.cell(row_idx, 0)
            text_cell = table.cell(row_idx, 1)
            
            title_cell.text = ad.get('title', '')
            text_cell.text = ad.get('text', '')
            
            # Белая заливка для ячеек данных
            for cell in [title_cell, text_cell]:
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                # Форматирование текста
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.size = font_size
        
        # Более агрессивное масштабирование шрифта в зависимости от количества строк
        initial_font_size = Pt(12)
        
        # Жесткое масштабирование шрифта по количеству строк
        if num_rows <= 5:
            font_size = Pt(12)  # 12pt для маленьких таблиц (до 5 строк)
        elif num_rows <= 7:
            font_size = Pt(10)  # 10pt для средних таблиц (6-7 строк)
        else:
            font_size = Pt(7)   # 7pt для больших таблиц (8+ строк)
        
        print(f"  📏 Установлен размер шрифта {font_size.pt}pt для таблицы с {num_rows} строками")
        
        # Применяем финальный размер шрифта ко всем ячейкам, включая заголовок
        for row_idx in range(num_rows):
            for col_idx in range(num_cols):
                cell = table.cell(row_idx, col_idx)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = font_size  # Устанавливаем размер для всего параграфа
                    for run in paragraph.runs:
                        run.font.size = font_size
                        if row_idx == 0:  # Для заголовка оставляем жирный шрифт
                            run.font.bold = True
        
        print(f"  ✅ Таблица создана с шрифтом {font_size.pt}pt")
    
    def _add_links_textbox(self, slide, unique_ads: List[Dict], left, top, width, height):
        """Добавить текстовый блок с уникальными ссылками под таблицей"""
        # Собираем уникальные ссылки
        unique_links = set()
        for ad in unique_ads:
            href = ad.get('href', '')
            if href:
                unique_links.add(href)
        
        print(f"  🔗 Найдено уникальных ссылок: {len(unique_links)}")
        
        # Формируем текст со ссылками
        links_text = " ".join(sorted(unique_links))
        
        # Создаем текстовый блок
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.clear()  # Очищаем стандартный параграф
        
        # Добавляем фиксированный текст (первая строка)
        p1 = text_frame.paragraphs[0]
        p1.text = "Ссылки на посадочные страницы:"
        p1.alignment = PP_ALIGN.LEFT
        
        for run in p1.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(7)
            run.font.color.rgb = RGBColor(102, 102, 102)  # #666666
        
        # Добавляем ссылки (вторая строка)
        if links_text:
            p2 = text_frame.add_paragraph()
            p2.text = links_text
            p2.alignment = PP_ALIGN.LEFT
            
            for run in p2.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(8)
                run.font.color.rgb = RGBColor(0, 0, 0)  # Черный
        
        print(f"  ✅ Добавлен блок со ссылками")
    
    def process_report(self, report: Dict) -> None:
        """Обработать один отчет"""
        print(f"\n{'='*60}")
        print(f"ОБРАБОТКА ОТЧЕТА #{report['id']}")
        print(f"{'='*60}")
        print(f"ID отчета: {report['id']}")
        print(f"ID договора: {report['id_contracts']}")
        print(f"ID заявки: {report['id_requests']}")
        print(f"Номер договора: {report['number_contract']}")
        print(f"Предмет договора: {report['subject_contract']}")
        
        # Загружаем данные из MinIO
        print(f"\n📥 Загрузка данных из MinIO...")
        
        # Формируем имена файлов с учетом номера отчета
        report_num = report['id']  # Используем id отчета как номер отчета
        ads_filename = f'ads_report_{report_num}.json'
        image_hashes_filename = f'image_hashes_report_{report_num}.json'
        
        # Загружаем файлы
        campaigns_data = self.load_file_from_minio(report['id'], 'campaigns.json')
        if not campaigns_data:
            print("❌ Не удалось загрузить campaigns.json")
            return
        
        ads_data = self.load_file_from_minio(report['id'], ads_filename)
        if not ads_data:
            print(f"❌ Не удалось загрузить {ads_filename}")
            return
        
        image_hashes_data = self.load_file_from_minio(report['id'], image_hashes_filename)
        if not image_hashes_data:
            print(f"❌ Не удалось загрузить {image_hashes_filename}")
            return
        
        # Фильтруем РСЯ-кампании
        print(f"\n🔍 Фильтрация РСЯ-кампаний...")
        rsy_campaigns = self.filter_rsy_campaigns(campaigns_data, report['campaign_ids'])
        
        if not rsy_campaigns:
            print("⚠ Не найдено РСЯ-кампаний для обработки")
            return
        
        # Создаем презентацию
        # Генерируем имя файла с датой и временем
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"presentation_{timestamp}.pptx"
        output_path = os.path.join(self.output_folder, filename)
        
        print(f"\n📊 Создание презентации...")
        success = self.create_presentation(rsy_campaigns, ads_data, image_hashes_data, output_path)
        
        if success:
            # читаем файл и переводим в байты
            file = io.BytesIO(open(output_path, 'rb').read())

            # # отправляем в S3
            # s3_file_path = os.getenv('S3_REPORT_PATH')
            filename = f"{report['id']}/" + filename
            # s3_file_path = '/'.join((s3_file_path, filename))
            # self.minio_client.put_object(
            #     self.bucket_name, s3_file_path, file, len(file.getvalue()))
            #
            # # записываем адрес (S3) в БД
            # write_s3path_to_bd(report.get('id'), os.getenv('PRESENTATION_COL_NAME'), s3_file_path)

            # удаляем файл
            os.remove(output_path)



            print(f"✅ Презентация успешно создана.")
            return file, filename
        else:
            print(f"❌ Не удалось создать презентацию")
    
    def run(self):
        """Основной метод запуска обработки"""
        print("🚀 Запуск генератора презентаций...")
        
        # Получаем отчеты для обработки
        reports = self.get_pending_reports()
        
        if not reports:
            print("📭 Нет отчетов со статусом 1 для обработки")
            return
        
        print(f"📋 Найдено отчетов для обработки: {len(reports)}")
        
        # Обрабатываем каждый отчет
        for report in reports:
            try:
                self.process_report(report)

            except Exception as e:
                print(f"❌ Ошибка при обработке отчета {report['id']}: {e}")
                import traceback
                traceback.print_exc()
                continue
        
        print(f"\n✅ Обработка завершена. Обработано отчетов: {len(reports)}")


def generate_presentation(report_id):
    """Главная функция"""
    try:
        generator = PresentationGenerator()
        from utils.postprocessing_report_file import get_report_by_id
        report = get_report_by_id(report_id)

        file, filename = generator.process_report(report)
        return file, filename

        # generator.run()
    except Exception as e:
        print(f"❌ Критическая ошибка: {e}")
        import traceback
        traceback.print_exc()



if __name__ == "__main__":
    print(generate_presentation(16))

